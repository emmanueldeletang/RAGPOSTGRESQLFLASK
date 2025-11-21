import os
import json
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

class DocumentProcessor:
    """Process different document types and extract text"""
    
    @staticmethod
    def process_pdf(file_path):
        """Extract text from PDF"""
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    
    @staticmethod
    def process_docx(file_path):
        """Extract text from DOCX and DOC (requires proper DOC support)"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            # Fallback: DOC files may not work with python-docx
            raise ValueError(f"Error processing Word document. Note: Legacy DOC format may not be fully supported. Consider converting to DOCX. Error: {e}")
    
    @staticmethod
    def process_pptx(file_path):
        """Extract text from PPTX and PPT (requires proper PPT support)"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            # Fallback: PPT files may not work with python-pptx
            raise ValueError(f"Error processing PowerPoint. Note: Legacy PPT format may not be fully supported. Consider converting to PPTX. Error: {e}")
    
    @staticmethod
    def process_xlsx(file_path):
        """Extract text from XLSX and XLS"""
        try:
            # pandas can handle both XLS and XLSX with openpyxl/xlrd
            df = pd.read_excel(file_path, sheet_name=None, engine=None)
            text = ""
            for sheet_name, sheet_df in df.items():
                text += f"Sheet: {sheet_name}\n"
                text += sheet_df.to_string(index=False) + "\n\n"
            return text
        except Exception as e:
            raise ValueError(f"Error processing Excel file. Ensure the file is not corrupted. Error: {e}")
    
    @staticmethod
    def process_csv(file_path):
        """Extract text from CSV"""
        df = pd.read_csv(file_path)
        return df.to_string(index=False)
    
    @staticmethod
    def process_json(file_path):
        """Extract text from JSON"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return json.dumps(data, indent=2)
    
    @staticmethod
    def process_txt(file_path):
        """Extract text from TXT"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    @staticmethod
    def process_file(file_path, file_type):
        """Process file based on type and extract text"""
        file_type = file_type.lower()
        
        processors = {
            'pdf': DocumentProcessor.process_pdf,
            'docx': DocumentProcessor.process_docx,
            'doc': DocumentProcessor.process_docx,
            'pptx': DocumentProcessor.process_pptx,
            'ppt': DocumentProcessor.process_pptx,
            'xlsx': DocumentProcessor.process_xlsx,
            'xls': DocumentProcessor.process_xlsx,
            'csv': DocumentProcessor.process_csv,
            'json': DocumentProcessor.process_json,
            'txt': DocumentProcessor.process_txt,
        }
        
        processor = processors.get(file_type)
        if processor:
            return processor(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
